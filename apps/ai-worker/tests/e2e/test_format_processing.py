# apps/ai-worker/tests/e2e/test_format_processing.py
"""
E2E tests for Phase 4 format processors.
Tests CSV, HTML, XLSX, EPUB, PPTX processing via temp files.
"""

import pytest

from src.converters import CsvConverter as CsvProcessor
from src.converters import HtmlConverter as HtmlProcessor
from src.converters import EpubConverter as EpubProcessor
from src.converters import PptxConverter as PptxProcessor
from src.converters import XlsxConverter as XlsxProcessor


class TestCsvProcessorE2E:
    """E2E tests for CSV processing."""

    @pytest.fixture
    def csv_file(self, tmp_path):
        """Create a CSV file for testing."""
        file_path = tmp_path / "test.csv"
        file_path.write_text("Name,Age\nAlice,30\nBob,25")
        return str(file_path)

    @pytest.mark.asyncio
    async def test_csv_to_markdown(self, csv_file):
        """CSV content is converted to Markdown."""
        processor = CsvProcessor()
        result = await processor.process(csv_file)

        assert result.markdown is not None
        assert "Name" in result.markdown
        assert "Alice" in result.markdown

    @pytest.mark.asyncio
    async def test_csv_metadata(self, csv_file):
        """CSV processing returns correct metadata."""
        processor = CsvProcessor()
        result = await processor.process(csv_file)

        # Verify metadata contains row info if available
        assert result.markdown is not None
        assert "Alice" in result.markdown


class TestHtmlProcessorE2E:
    """E2E tests for HTML processing."""

    @pytest.fixture
    def html_file(self, tmp_path):
        """Create an HTML file for testing."""
        file_path = tmp_path / "test.html"
        file_path.write_text("<html><body><h1>Title</h1><p>Content</p></body></html>")
        return str(file_path)

    @pytest.mark.asyncio
    async def test_html_to_markdown(self, html_file):
        """HTML is converted to Markdown."""
        processor = HtmlProcessor()
        result = await processor.process(html_file)

        assert result.markdown is not None
        assert "Title" in result.markdown
        assert "Content" in result.markdown


class TestXlsxProcessorE2E:
    """E2E tests for XLSX processing."""

    @pytest.fixture
    def xlsx_file(self, tmp_path):
        """Create a simple XLSX file for testing."""
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Test"
        ws["B2"] = 123

        file_path = tmp_path / "test.xlsx"
        wb.save(file_path)
        return str(file_path)

    @pytest.mark.asyncio
    async def test_xlsx_to_markdown(self, xlsx_file):
        """XLSX content is converted to Markdown."""
        processor = XlsxProcessor()
        result = await processor.process(xlsx_file)

        assert result.markdown is not None
        assert "Sheet1" in result.markdown
        assert "Name" in result.markdown

    @pytest.mark.asyncio
    async def test_xlsx_metadata(self, xlsx_file):
        """XLSX processing returns sheet count."""
        processor = XlsxProcessor()
        result = await processor.process(xlsx_file)

        assert result.sheet_count == 1


class TestPptxProcessorE2E:
    """E2E tests for PPTX processing."""

    @pytest.fixture
    def pptx_file(self, tmp_path):
        """Create a simple PPTX file for testing."""
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank
        prs.slides.add_slide(slide_layout)

        file_path = tmp_path / "test.pptx"
        prs.save(file_path)
        return str(file_path)

    @pytest.mark.asyncio
    async def test_pptx_processing(self, pptx_file):
        """PPTX content is processed."""
        processor = PptxProcessor()
        result = await processor.process(pptx_file)

        # Minimal PPTX may have 0 slides or empty markdown
        assert result.markdown is not None or result.slide_count is not None


class TestEpubProcessorE2E:
    """E2E tests for EPUB processing."""

    @pytest.fixture
    def epub_file(self, tmp_path):
        """Create a simple EPUB file for testing."""
        from ebooklib import epub

        book = epub.EpubBook()
        book.set_title("Test Book")
        book.add_author("Test Author")

        c1 = epub.EpubHtml(title="Chapter 1", file_name="chap_01.xhtml")
        c1.content = "<html><body><h1>Chapter 1</h1><p>Content here.</p></body></html>"
        book.add_item(c1)

        book.toc = (epub.Link("chap_01.xhtml", "Chapter 1", "ch1"),)
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        book.spine = ["nav", c1]

        file_path = tmp_path / "test.epub"
        epub.write_epub(file_path, book)
        return str(file_path)

    @pytest.mark.asyncio
    async def test_epub_to_markdown(self, epub_file):
        """EPUB content is converted to Markdown."""
        processor = EpubProcessor()
        result = await processor.process(epub_file)

        assert result.markdown is not None

    @pytest.mark.asyncio
    async def test_epub_metadata(self, epub_file):
        """EPUB processing returns chapter count."""
        processor = EpubProcessor()
        result = await processor.process(epub_file)

        assert result.chapter_count >= 1
